#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Flask backend for music teaching platform.
Combines audio analysis (proxy.py) and PPT generation (ppt_generator.py).
"""
from flask import Flask, send_file, request, jsonify
from flask_cors import CORS
from pptx import Presentation
from pptx.util import Inches, Pt
import requests
import io
import json
import os

BASE_DIR = os.path.abspath(os.path.dirname(__file__))
DATA_DIR = os.path.join(BASE_DIR, 'songs')

app = Flask(__name__)
CORS(app)

# Audio analysis API configuration
AUDD_API_KEY = os.environ.get('AUDD_API_KEY', 'a614eaee55a298b2625726b1915f3a0a')
AUDD_API_URL = 'https://api.audd.io/analysis/'


# ============ Audio Analysis Routes ============

@app.route('/api/analyze', methods=['POST'])
def analyze():
    """Analyze uploaded audio file using AudD API."""
    if 'file' not in request.files:
        return jsonify({'error': 'No file uploaded'}), 400

    file = request.files['file']
    files = {'file': (file.filename, file.stream, file.mimetype)}
    data = {
        'api_token': AUDD_API_KEY,
        'return': 'genre, mood'
    }

    try:
        resp = requests.post(AUDD_API_URL, files=files, data=data)
        return jsonify(resp.json())
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/scores', methods=['GET'])
def get_scores():
    """List available score files."""
    scores_dir = os.path.join(BASE_DIR, '..', 'frontend', 'public', 'scores')
    try:
        if not os.path.exists(scores_dir):
            return jsonify({'scores': []})
        files = os.listdir(scores_dir)
        files = [f for f in files if os.path.isfile(os.path.join(scores_dir, f))]
        return jsonify({'scores': files})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


# ============ PPT Generation Routes ============

def create_ppt_from_metadata(meta):
    """Create a PowerPoint presentation from song metadata."""
    prs = Presentation()

    # Title slide
    try:
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        if 'title' in meta:
            slide.shapes.title.text = meta.get('title', '')
        if 'composer' in meta:
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = meta.get('composer', '')
    except Exception:
        pass

    # Score images
    for img in meta.get('scoreImages', []):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        left = top = Inches(0.5)
        width = prs.slide_width - Inches(1)
        img_path = img
        if not os.path.isabs(img_path):
            img_path = os.path.join(BASE_DIR, img_path)
        try:
            slide.shapes.add_picture(img_path, left, top, width=width)
        except Exception:
            tx_box = slide.shapes.add_textbox(left, top, width, Inches(1))
            tf = tx_box.text_frame
            tf.text = f'无法加载图片: {img}'

    # Practice tips
    if meta.get('practiceTips'):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = '练习要点'
        body = None
        if len(slide.placeholders) > 1:
            body = slide.placeholders[1].text_frame
            body.clear()
        else:
            left = Inches(0.5)
            top = Inches(1.5)
            width = prs.slide_width - Inches(1)
            body = slide.shapes.add_textbox(left, top, width, Inches(3)).text_frame

        for tip in meta['practiceTips']:
            p = body.add_paragraph()
            p.text = '• ' + tip
            p.level = 0
            p.font.size = Pt(20)

    # Teacher notes
    if meta.get('teacherNotes'):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = '教师提示'
        if len(slide.placeholders) > 1:
            tf = slide.placeholders[1].text_frame
            tf.text = meta.get('teacherNotes')
        else:
            left = Inches(0.5)
            top = Inches(1.5)
            width = prs.slide_width - Inches(1)
            slide.shapes.add_textbox(left, top, width, Inches(3)).text_frame.text = meta.get('teacherNotes')

    # Audio info
    if meta.get('audio'):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = '伴奏 / 音频'
        if len(slide.placeholders) > 1:
            tf = slide.placeholders[1].text_frame
            tf.text = '伴奏文件：'
            p = tf.add_paragraph()
            p.text = meta.get('audio')
        else:
            left = Inches(0.5)
            top = Inches(1.5)
            width = prs.slide_width - Inches(1)
            tb = slide.shapes.add_textbox(left, top, width, Inches(1)).text_frame
            tb.text = '伴奏文件：'
            p = tb.add_paragraph()
            p.text = meta.get('audio')

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


@app.route('/api/generate_ppt', methods=['GET'])
def generate_ppt():
    """Generate a PowerPoint presentation for a song."""
    song_id = request.args.get('song')
    if not song_id:
        return jsonify({'error': 'missing song id'}), 400

    meta_path = os.path.join(DATA_DIR, f'{song_id}.json')
    if not os.path.exists(meta_path):
        return jsonify({'error': 'metadata not found', 'path': meta_path}), 404

    with open(meta_path, 'r', encoding='utf-8') as f:
        meta = json.load(f)

    try:
        pptx_io = create_ppt_from_metadata(meta)
        filename = f"{meta.get('title', song_id)}.pptx"
        return send_file(
            pptx_io,
            as_attachment=True,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            download_name=filename
        )
    except Exception as e:
        return jsonify({'error': '生成PPT失败', 'message': str(e)}), 500


# ============ Health Check ============

@app.route('/api/health', methods=['GET'])
def health():
    """Health check endpoint."""
    return jsonify({'status': 'ok'})


if __name__ == '__main__':
    app.run(host='0.0.0.0', port=5000, debug=True)
