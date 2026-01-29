"""
PPT generation service.
"""
import io
import json
import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from app.schemas.ppt import PPTMetadata
from app.core.exceptions import NotFoundError

# Path to song metadata
SONGS_DIR = Path(__file__).parent.parent.parent / "songs"


def load_song_metadata(song_id: str) -> dict:
    """Load song metadata from JSON file."""
    meta_path = SONGS_DIR / f"{song_id}.json"
    if not meta_path.exists():
        raise NotFoundError(f"Song metadata '{song_id}' not found", field="song")

    with open(meta_path, "r", encoding="utf-8") as f:
        return json.load(f)


def create_ppt_from_metadata(meta: dict) -> io.BytesIO:
    """
    Create a PowerPoint presentation from song metadata.

    Args:
        meta: Song metadata dictionary

    Returns:
        BytesIO buffer containing the PPTX file
    """
    prs = Presentation()

    # Title slide
    try:
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        if "title" in meta:
            slide.shapes.title.text = meta.get("title", "")
        if "composer" in meta:
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = meta.get("composer", "")
    except Exception:
        pass

    # Score images
    for img in meta.get("scoreImages", []):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        left = top = Inches(0.5)
        width = prs.slide_width - Inches(1)
        img_path = img
        if not os.path.isabs(img_path):
            img_path = str(SONGS_DIR.parent / img_path)
        try:
            slide.shapes.add_picture(img_path, left, top, width=width)
        except Exception:
            tx_box = slide.shapes.add_textbox(left, top, width, Inches(1))
            tf = tx_box.text_frame
            tf.text = f"无法加载图片: {img}"

    # Practice tips
    if meta.get("practiceTips"):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "练习要点"
        body = None
        if len(slide.placeholders) > 1:
            body = slide.placeholders[1].text_frame
            body.clear()
        else:
            left = Inches(0.5)
            top = Inches(1.5)
            width = prs.slide_width - Inches(1)
            body = slide.shapes.add_textbox(left, top, width, Inches(3)).text_frame

        for tip in meta["practiceTips"]:
            p = body.add_paragraph()
            p.text = "• " + tip
            p.level = 0
            p.font.size = Pt(20)

    # Teacher notes
    if meta.get("teacherNotes"):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "教师提示"
        if len(slide.placeholders) > 1:
            tf = slide.placeholders[1].text_frame
            tf.text = meta.get("teacherNotes")
        else:
            left = Inches(0.5)
            top = Inches(1.5)
            width = prs.slide_width - Inches(1)
            slide.shapes.add_textbox(left, top, width, Inches(3)).text_frame.text = (
                meta.get("teacherNotes")
            )

    # Audio info
    if meta.get("audio"):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "伴奏 / 音频"
        if len(slide.placeholders) > 1:
            tf = slide.placeholders[1].text_frame
            tf.text = "伴奏文件："
            p = tf.add_paragraph()
            p.text = meta.get("audio")
        else:
            left = Inches(0.5)
            top = Inches(1.5)
            width = prs.slide_width - Inches(1)
            tb = slide.shapes.add_textbox(left, top, width, Inches(1)).text_frame
            tb.text = "伴奏文件："
            p = tb.add_paragraph()
            p.text = meta.get("audio")

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


def generate_ppt(song_id: str) -> tuple[io.BytesIO, str]:
    """
    Generate PPT for a song.

    Args:
        song_id: Song identifier

    Returns:
        Tuple of (BytesIO buffer, filename)
    """
    meta = load_song_metadata(song_id)
    pptx_io = create_ppt_from_metadata(meta)
    filename = f"{meta.get('title', song_id)}.pptx"
    return pptx_io, filename
