#!/usr/bin/env python3
# -*- coding: utf-8 -*-
from flask import Flask, send_file, request, jsonify
from pptx import Presentation
from pptx.util import Inches, Pt
import io, json, os

BASE_DIR = os.path.abspath(os.path.dirname(__file__))
DATA_DIR = os.path.join(BASE_DIR, 'songs')

app = Flask(__name__)

def create_ppt_from_metadata(meta):
    prs = Presentation()
    # Title slide
    try:
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        if 'title' in meta:
            slide.shapes.title.text = meta.get('title', '')
        if 'composer' in meta:
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = meta.get('composer', '')
    except Exception:
        pass

    # Score images
    for img in meta.get('scoreImages', []):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        left = top = Inches(0.5)
        width = prs.slide_width - Inches(1)
        # Resolve path relative to project
        img_path = img
        if not os.path.isabs(img_path):
            img_path = os.path.join(BASE_DIR, img_path)
        try:
            slide.shapes.add_picture(img_path, left, top, width=width)
        except Exception:
            # Fallback: put a text box indicating missing image
            tx_box = slide.shapes.add_textbox(left, top, width, Inches(1))
            tf = tx_box.text_frame
            tf.text = f'无法加载图片: {img}'

    # Practice tips
    if meta.get('practiceTips'):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = '练习要点'
        # Use the first placeholder for body
        body = None
        if len(slide.placeholders) > 1:
            body = slide.placeholders[1].text_frame
            body.clear()
        else:
            left = Inches(0.5)
            top = Inches(1.5)
            width = prs.slide_width - Inches(1)
            body = slide.shapes.add_textbox(left, top, width, Inches(3)).text_frame

        for tip in meta['practiceTips']:
            p = body.add_paragraph()
            p.text = u'• ' + tip
            p.level = 0
            p.font.size = Pt(20)

    # Teacher notes
    if meta.get('teacherNotes'):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = '教师提示'
        if len(slide.placeholders) > 1:
            tf = slide.placeholders[1].text_frame
            tf.text = meta.get('teacherNotes')
        else:
            left = Inches(0.5)
            top = Inches(1.5)
            width = prs.slide_width - Inches(1)
            slide.shapes.add_textbox(left, top, width, Inches(3)).text_frame.text = meta.get('teacherNotes')

    # Audio info (link)
    if meta.get('audio'):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = '伴奏 / 音频'
        if len(slide.placeholders) > 1:
            tf = slide.placeholders[1].text_frame
            tf.text = '伴奏文件：'
            p = tf.add_paragraph()
            p.text = meta.get('audio')
        else:
            left = Inches(0.5)
            top = Inches(1.5)
            width = prs.slide_width - Inches(1)
            tb = slide.shapes.add_textbox(left, top, width, Inches(1)).text_frame
            tb.text = '伴奏文件：'
            p = tb.add_paragraph()
            p.text = meta.get('audio')

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


@app.route('/generate_ppt')
def generate_ppt():
    song_id = request.args.get('song')
    if not song_id:
        return jsonify({"error": "missing song id"}), 400
    meta_path = os.path.join(DATA_DIR, f'{song_id}.json')
    if not os.path.exists(meta_path):
        return jsonify({"error": "metadata not found", "path": meta_path}), 404
    with open(meta_path, 'r', encoding='utf-8') as f:
        meta = json.load(f)
    try:
        pptx_io = create_ppt_from_metadata(meta)
        filename = f"{meta.get('title', song_id)}.pptx"
        return send_file(pptx_io, as_attachment=True, mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation', download_name=filename)
    except Exception as e:
        return jsonify({"error": "生成PPT失败", "message": str(e)}), 500


if __name__ == '__main__':
    app.run(host='0.0.0.0', port=5000, debug=True)
